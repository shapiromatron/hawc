import codecs
from datetime import datetime
from io import StringIO
import logging
import os
import re
import subprocess
import tempfile
import urllib2

from django.conf import settings
from django.http.request import HttpRequest
from django.template import RequestContext
from django.template.loader import render_to_string

from pptx import Presentation
from pptx.util import Inches, Pt


logger = logging.getLogger(__name__)


def _get_css_styles():
    txt = u''

    path = os.path.join(
        settings.PROJECT_PATH, r'static/css/hawc_d3_aggregated.css')

    if os.path.exists(path):
        with codecs.open(path, 'r', 'UTF-8') as f:
            txt = unicode(f.read())
        txt = r'<defs style="hawc-styles"><style type="text/css"><![CDATA[{0}]]></style></defs>'.format(txt)  # noqa

    return txt


# load CSS into memory upon runtime (only do it once)
D3_CSS_STYLES = _get_css_styles()


class SVGConverter(object):

    def __init__(self, svg, url, width, height):
        self.url = url
        self.width = width
        self.height = height
        self.tempfns = []

        svg = svg.decode('base64')\
            .replace('%u', '\\u')\
            .decode('unicode_escape')
        self.svg = urllib2.unquote(svg)

    def to_svg(self):
        # add embedded styles
        svg = self.svg

        # add CSS styles
        match = re.search(r'<svg [^>]+>', svg)
        insertion_point = match.end()

        # Manually add our CSS styles from a file. Because there are problems
        # inserting CDATA using a python etree, we use a regex instead
        svg = (svg[:insertion_point] + D3_CSS_STYLES + svg[insertion_point:])

        return svg

    def to_png(self, delete_and_return_object=True):
        logger.info('Converting svg -> html -> png')
        content = None
        try:
            png = self.get_tempfile(suffix='.png')
            self._rasterize(png)
            if delete_and_return_object:
                with open(png, 'rb') as f:
                    content = f.read()
            else:
                content = png
        except Exception as e:
            logger.error(e.message, exc_info=True)
        finally:
            self.cleanup()
        return content

    def to_pdf(self):
        logger.info('Converting svg -> html -> pdf')
        content = None
        try:
            pdf = self.get_tempfile(suffix='.pdf')
            self._rasterize(pdf)
            with open(pdf, 'rb') as f:
                content = f.read()
        except Exception as e:
            logger.error(e.message, exc_info=True)
        finally:
            self.cleanup()
        return content

    def _pptx_add_title(self, slide):
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
        tf = txBox.textframe
        now = datetime.now().strftime("%B %d %Y, %I:%M %p")
        tf.text = "HAWC visualization generated on {}".format(now)
        tf.paragraphs[0].alignment = 2

    def _pptx_add_url(self, slide):
        txBox = slide.shapes.add_textbox(
            Inches(0), Inches(7.0), Inches(10), Inches(0.5))
        tf = txBox.textframe
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = self.url
        run.hyperlink.address = self.url
        p.font.size = Pt(12)
        p.alignment = 2

    def _pptx_add_png(self, slide, png_fn):
        # add picture, after getting proper dimensions
        max_height = 6.
        max_width = 9.
        left = Inches(0.5)
        top = Inches(0.75)
        width_to_height_ratio = float(self.width)/self.height
        if width_to_height_ratio > (max_width/max_height):
            width = Inches(max_width)
            height = width/width_to_height_ratio
            top = top + int((Inches(max_height)-height)*0.5)
        else:
            height = Inches(max_height)
            width = height*width_to_height_ratio
            left = left + int((Inches(max_width)-width)*0.5)

        slide.shapes.add_picture(png_fn, left, top, width, height)

    def _pptx_add_hawc_logo(self, slide):
        left = Inches(8.55)
        top = Inches(6.65)
        width = Inches(1.4)
        height = Inches(0.8)
        logo_location = os.path.join(settings.STATIC_ROOT, "img/HAWC-120.png")
        slide.shapes.add_picture(logo_location, left, top, width, height)

    def to_pptx(self):
        logger.info('Converting svg -> html -> png -> pptx')
        content = None
        try:
            # convert to png
            png_fn = self.to_png(delete_and_return_object=False)

            # create blank presentation slide layout
            prs = Presentation()
            blank_slidelayout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slidelayout)

            self._pptx_add_title(slide)
            self._pptx_add_url(slide)
            self._pptx_add_png(slide, png_fn)
            self._pptx_add_hawc_logo(slide)

            # save as object
            content = StringIO()
            prs.save(content)
            content.seek(0)

        except Exception as e:
            logger.error(e.message, exc_info=True)
        finally:
            self.cleanup()
        return content

    def get_tempfile(self, prefix='hawc-', suffix='.txt'):
        _, fn = tempfile.mkstemp(prefix=prefix, suffix=suffix)
        self.tempfns.append(fn)
        return fn

    def cleanup(self):
        for fn in self.tempfns:
            os.remove(fn)

    def _to_html(self):
        # return rendered html absolute filepath
        request = HttpRequest()
        context = RequestContext(request, dict(
            svg=self.svg,
            css=D3_CSS_STYLES
        ))
        html = render_to_string('rasterize.html', context).encode('UTF-8')
        fn = self.get_tempfile(suffix='.html')
        with open(fn, 'wb') as f:
            f.write(html)
        return fn

    def _rasterize(self, out_fn):
        phantom = settings.PHANTOMJS_PATH
        rasterize = os.path.join(
            settings.PROJECT_PATH, 'static/js/rasterize.js')
        html_fn = self._to_html()
        try:
            commands = [phantom, rasterize, html_fn, out_fn]
            subprocess.call(commands)
            logger.info('Conversion successful')
        except Exception as e:
            logger.error(e.message, exc_info=True)
