from __future__ import absolute_import
from datetime import datetime
import codecs
import os
import random
import re
import string
import subprocess
import tempfile
from StringIO import StringIO
import urllib2
import xml.etree.ElementTree as ET

from django.conf import settings
from django.core.cache import cache
from django.http.request import HttpRequest
from django.template import RequestContext
from django.template.loader import render_to_string

from celery import shared_task
from celery.utils.log import get_task_logger
import requests
from pptx import Presentation
from pptx.util import Inches, Pt


logger = get_task_logger(__name__)


# load CSS into memory upon runtime (only do it once)
D3_CSS_STYLES = u""
D3_CSS_PATH = os.path.abspath(
    os.path.join(settings.PROJECT_PATH,
                 r'static/css/hawc_d3_aggregated.css'))
if os.path.exists(D3_CSS_PATH):
    with codecs.open(D3_CSS_PATH, 'r', 'UTF-8') as f:
        D3_CSS_STYLES = unicode(f.read())


class SVGConverter():

    def __init__(self, svg, url=None, width=None, height=None):
        self.url = url
        self.width = width
        self.height = height
        self.tempfns = []

        svg = svg.decode('base64')\
            .replace('%u', '\\u')\
            .decode('unicode_escape')
        self.svg = urllib2.unquote(svg)

    def get_svg_with_embedded_styles(self):
        svg = self.svg

        # add CSS styles
        match = re.search(ur'<svg [^>]+>', svg)
        insertion_point = match.end()

        # Manually add our CSS styles from a file. Because there are problems
        # inserting CDATA using a python etree, we use a regex instead
        styles = r'<defs style="hawc-styles"><style type="text/css"><![CDATA[{0}]]></style></defs>'.format(D3_CSS_STYLES)
        svg = (svg[:insertion_point] + styles + svg[insertion_point:])

        return svg

    def get_tempfile(self, prefix='hawc-', suffix='.txt'):
        _, fn = tempfile.mkstemp(prefix=prefix, suffix=suffix)
        self.tempfns.append(fn)
        return fn

    def cleanup(self):
        for fn in self.tempfns:
            os.remove(fn)

    def convert_to_html(self):
        request = HttpRequest()
        context = RequestContext(request, dict(
            svg=self.svg,
            css=D3_CSS_STYLES
        ))
        html = render_to_string('rasterize.html', context).encode('UTF-8')
        fn = self.get_tempfile(suffix='.html')
        with open(fn, 'wb') as f:
            f.write(html)
        return fn

    @shared_task
    def convert_to_png(self, delete_and_return_object=True):
        logger.info('Converting svg -> html -> png')
        png = self.get_tempfile(suffix='.png')
        self.rasterize(png)
        content = None
        if delete_and_return_object:
            with open(png, 'rb') as f:
                content = f.read()
            self.cleanup()
            return content
        else:
            self.png_fn = png

    @shared_task
    def convert_to_pdf(self):
        logger.info('Converting svg -> html -> pdf')
        pdf = self.get_tempfile(suffix='.pdf')
        self.rasterize(pdf)
        content = None
        with open(pdf, 'rb') as f:
            content = f.read()
        self.cleanup()
        return content

    def rasterize(self, out_fn):
        phantom = settings.PHANTOMJS_PATH
        rasterize = os.path.join(settings.PROJECT_PATH, 'static', 'js', 'rasterize.js')
        html_fn = self.convert_to_html()
        try:
            commands = [phantom, rasterize, html_fn, out_fn]
            subprocess.call(commands)
            logger.info('Conversion successful')
        except Exception as e:
            logger.error(e.message, exc_info=True)

    @shared_task
    def convert_to_pptx(_, self):
        logger.info('Converting svg to pptx')
        pptx = None

        try:
            # create blank presentation slide layout
            prs = Presentation()
            blank_slidelayout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slidelayout)

            # add title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
            tf = txBox.textframe
            now = datetime.now().strftime("%B %d %Y, %I:%M %p")
            tf.text = "HAWC visualization generated on {now}".format(now=now)
            tf.paragraphs[0].alignment = 2

            # add website address
            txBox = slide.shapes.add_textbox(Inches(0), Inches(7.0), Inches(10), Inches(0.5))
            tf = txBox.textframe
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = self.url
            run.hyperlink.address = self.url
            p.font.size = Pt(12)
            p.alignment = 2

            # add picture, after getting proper dimensions
            max_height = 6.
            max_width = 9.
            left = Inches(0.5)
            top = Inches(0.75)
            width_to_height_ratio = float(self.width)/self.height
            if width_to_height_ratio > (max_width/max_height):
                width = Inches(max_width)
                height = width/width_to_height_ratio
                top = top + int((Inches(max_height)-height)*0.5)
            else:
                height = Inches(max_height)
                width = height*width_to_height_ratio
                left = left + int((Inches(max_width)-width)*0.5)

            slide.shapes.add_picture(self.png_fn, left, top, width, height)

            # add HAWC logo
            left = Inches(8.55)
            top = Inches(6.65)
            width = Inches(1.4)
            height = Inches(0.8)
            logo_location = os.path.join(settings.STATIC_ROOT, "img/HAWC-120.png")
            slide.shapes.add_picture(logo_location, left, top, width, height)

            # save as object
            output = StringIO()
            prs.save(output)
            output.seek(0)

            pptx = output
            logger.info('Conversion successful')

        except Exception as e:
            logger.error(e.message, exc_info=True)

        self.cleanup()
        return pptx


@shared_task
def get_chemspider_details(cas_number):

    cache_name = 'chemspider-{cas_number}'.format(cas_number=cas_number.replace(' ', '_'))
    d = cache.get(cache_name)
    if d is None:
        d = {"status": "failure"}
        if cas_number:
            try:
                # get chemspider chem id
                url = 'http://www.chemspider.com/Search.asmx/SimpleSearch'
                payload = {'query': cas_number,
                           'token': settings.CHEMSPIDER_TOKEN}
                r = requests.post(url, data=payload)

                id_val = re.search(r'\<int\>(\d+)\</int\>', r.content).group(1)

                # get details
                url = 'http://www.chemspider.com/MassSpecAPI.asmx/GetExtendedCompoundInfo'
                payload = {'CSID': id_val,
                           'token': settings.CHEMSPIDER_TOKEN}
                r = requests.post(url, data=payload)
                xml = ET.fromstring(r.content)
                namespace = '{http://www.chemspider.com/}'
                d['CommonName'] = xml.find('{ns}CommonName'.format(ns=namespace)).text
                d['SMILES'] = xml.find('{ns}SMILES'.format(ns=namespace)).text
                d['MW'] = xml.find('{ns}MolecularWeight'.format(ns=namespace)).text

                # get image
                url = 'http://www.chemspider.com/Search.asmx/GetCompoundThumbnail'
                payload = {'id': id_val,
                           'token': settings.CHEMSPIDER_TOKEN}
                r = requests.post(url, data=payload)
                xml = ET.fromstring(r.content)
                d['image'] = xml.text

                # call it a success if we made it here
                d['status'] = 'success'

                logger.info('setting cache: {cache_name}'.format(cache_name=cache_name))
                cache.set(cache_name, d)
            except Exception as e:
                logger.error(e.message, exc_info=True)

    return d
