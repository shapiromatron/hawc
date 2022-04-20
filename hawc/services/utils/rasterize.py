import base64
import binascii
import logging
import os
import re
import shlex
import subprocess
import tempfile
from datetime import datetime
from io import BytesIO
from typing import Optional
from urllib import parse

from django.conf import settings
from django.template.loader import render_to_string
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


class HawcStyles:
    @property
    def for_svg(self):
        # lazy evaluation
        style = getattr(self, "_for_svg", None)
        if style is None:
            style = f'<defs><style type="text/css"><![CDATA[{self._get_styles()}]]></style></defs>'
            setattr(self, "_for_svg", style)
        return style

    @property
    def for_html(self):
        # lazy evaluation
        style = getattr(self, "_for_html", None)
        if style is None:
            style = f'<style type="text/css">{self._get_styles()}</style>'
            setattr(self, "_for_html", style)
        return style

    def _get_styles(self):
        def remove_spacing(txt, character):
            txt = re.sub(character + " ", character, txt)
            return re.sub(" " + character, character, txt)

        def get_styles(fn):
            # Only using D3.css as adding other style-sheets with non SVG
            # styles can potentially break the SVG processor.
            texts = []
            with open(fn, "r") as f:
                texts.append(f.read())
            return " ".join(texts)

        def compact(txt):
            txt = re.sub(r"\/\*[^(\*\/)]+\*\/", " ", txt)
            txt = re.sub(r"\n", " ", txt)
            txt = re.sub(r" >", " ", txt)  # can break illustrator
            txt = re.sub(r"[ ]+", " ", txt)
            txt = remove_spacing(txt, ":")
            txt = remove_spacing(txt, "{")
            txt = remove_spacing(txt, "}")
            txt = remove_spacing(txt, ";")
            txt = re.sub(r"}", r"}\n", txt)
            return txt

        css_file = str(settings.PROJECT_PATH / "static/css/d3.css")
        return compact(get_styles(css_file))


Styles = HawcStyles()


RE_SVG = re.compile(r"<svg.*?>.+?</svg>", re.DOTALL)


class SVGConverter:
    def __init__(self, svg: str, url: str, width: int, height: int):
        self.svg = svg
        self.url = url
        self.width = width
        self.height = height

    @classmethod
    def decode_svg(cls, svg) -> str:
        """Decode base64 encoded svg data

        Args:
            svg: base64 encoded svg

        Raises:
            ValueError: If data cannot be decoded or is invalid SVG
        """
        try:
            decoded = (
                base64.b64decode(svg)
                .decode("utf8")
                .replace("%u", "\\u")
                .encode()
                .decode("unicode-escape")
            )
        except (binascii.Error, ValueError, UnicodeDecodeError):
            raise ValueError("Invalid base64 encoding")

        # ensure svg-like
        unquoted = parse.unquote(decoded, encoding="ISO-8859-1")
        if not RE_SVG.fullmatch(unquoted):
            raise ValueError("Invalid SVG")

        return unquoted

    def to_svg(self):
        svg = self.svg

        # add CSS style definitions
        match = re.search(r"<svg [^>]+>", svg)
        insertion_point = match.end()

        # Manually add our CSS styles from a file. Because there are problems
        # inserting CDATA using a python etree, we use a regex instead
        return "".join([svg[:insertion_point], Styles.for_svg, svg[insertion_point:]])

    def to_png(self) -> Optional[bytes]:
        return self._rasterize(suffix=".png")

    def to_pdf(self) -> Optional[bytes]:
        return self._rasterize(suffix=".pdf")

    def to_pptx(self) -> Optional[bytes]:
        logger.info("Converting svg -> html -> png -> pptx")

        # convert to png
        png = self.to_png()
        if png is None:
            return None

        # create blank presentation slide layout
        pres = Presentation()
        blank_slidelayout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(blank_slidelayout)

        self._pptx_add_title(slide)
        self._pptx_add_url(slide)
        self._pptx_add_png(slide, BytesIO(png))
        self._pptx_add_hawc_logo(slide)

        # save as object
        content = BytesIO()
        pres.save(content)
        content.seek(0)

        return content.read()

    def _pptx_add_title(self, slide):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
        tf = txBox.text_frame
        now = datetime.now().strftime("%B %d %Y, %I:%M %p")
        tf.text = f"HAWC visualization generated on {now}"
        tf.paragraphs[0].alignment = 2

    def _pptx_add_url(self, slide):
        txBox = slide.shapes.add_textbox(Inches(0), Inches(7.0), Inches(10), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = self.url
        run.hyperlink.address = self.url
        p.font.size = Pt(12)
        p.alignment = 2

    def _pptx_add_png(self, slide, png_fn):
        # add picture, after getting proper dimensions
        max_height = 6.0
        max_width = 9.0
        left = Inches(0.5)
        top = Inches(0.75)
        width_to_height_ratio = float(self.width) / self.height
        if width_to_height_ratio > (max_width / max_height):
            width = Inches(max_width)
            height = width / width_to_height_ratio
            top = top + int((Inches(max_height) - height) * 0.5)
        else:
            height = Inches(max_height)
            width = height * width_to_height_ratio
            left = left + int((Inches(max_width) - width) * 0.5)

        slide.shapes.add_picture(png_fn, left, top, width, height)

    def _pptx_add_hawc_logo(self, slide):
        left = Inches(8.55)
        top = Inches(6.65)
        width = Inches(1.4)
        height = Inches(0.8)
        logo_location = os.path.join(settings.STATICFILES_DIRS[0], "img/HAWC-400.png")
        slide.shapes.add_picture(logo_location, left, top, width, height)

    def _to_html(self, f):
        # return rendered html absolute filepath
        context = dict(svg=self.svg, css=Styles.for_html)
        html = render_to_string("rasterize.html", context).encode("UTF-8")
        f.write(html)

    def _rasterize(self, suffix) -> Optional[bytes]:
        phantomjs_env = os.environ.copy()
        phantomjs_env.update(settings.PHANTOMJS_ENV)
        rasterize = str(settings.PROJECT_PATH / "static/js/rasterize.js")
        with (
            tempfile.NamedTemporaryFile("wb", prefix="hawc-", suffix=".html") as input,
            tempfile.NamedTemporaryFile("rb", prefix="hawc-", suffix=suffix) as output,
        ):
            self._to_html(input)
            try:
                commands = shlex.split(
                    " ".join([settings.PHANTOMJS_PATH, rasterize, input.name, output.name])
                )
                subprocess.run(commands, env=phantomjs_env, check=True)
                return output.read()
            except subprocess.CalledProcessError as err:
                logger.error(err, exc_info=True)
